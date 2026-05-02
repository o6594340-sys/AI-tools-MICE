import io
import json
import os

from anthropic import Anthropic
from dotenv import load_dotenv
from flask import Flask, Response, jsonify, render_template, request, stream_with_context

import prompts as p

load_dotenv()

app = Flask(__name__)
client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

MODEL = "claude-sonnet-4-6"
MAX_TOKENS = 2000


def build_user_message(mode, data):
    if mode == "brief":
        lang = data.get("language", "русский")
        return f"Вот бриф клиента:\n\n{data.get('brief_text', '')}\n\nПредпочтительный язык ответа: {lang}"

    if mode == "objection":
        obj_type = data.get("objection_type", "определить автоматически")
        lang = data.get("language", "русский")
        context = data.get("context", "").strip()
        msg = f"Тип возражения: {obj_type}\n\nТекст возражения клиента:\n{data.get('objection_text', '')}"
        if context:
            msg += f"\n\nКонтекст / направление: {context}"
        msg += f"\n\nЯзык ответа: {lang}"
        return msg

    if mode == "supplier":
        goals = data.get("goals", [])
        custom_goal = data.get("custom_goal", "").strip()
        if custom_goal:
            goals = goals + [custom_goal]
        goals_text = ", ".join(goals) if goals else "не указана"
        tone = data.get("tone", "нейтральный деловой")
        return (
            f"Контекст / текст от поставщика:\n{data.get('context', '')}\n\n"
            f"Цель письма: {goals_text}\n\n"
            f"Тон письма: {tone}"
        )

    if mode == "dmc":
        tone = data.get("tone", "сдержанный профессиональный")
        return f"Текст от DMC:\n\n{data.get('dmc_text', '')}\n\nТон вывода: {tone}"

    if mode == "dmc_refine":
        style = data.get("style", "более продающе")
        style_instruction = p.DMC_STYLE_MAP.get(style, "")
        return f"ИНСТРУКЦИЯ ПО СТИЛЮ:\n{style_instruction}\n\nТЕКСТ ДЛЯ ПЕРЕРАБОТКИ:\n{data.get('current_text', '')}"

    if mode == "followup":
        proposal = data.get("proposal", "").strip() or "не указано"
        days = data.get("days", "4–7")
        tone = data.get("tone", "деловой")
        lang = data.get("language", "русский")
        return (
            f"Что было предложено: {proposal}\n\n"
            f"Дней без ответа: {days}\n\n"
            f"Тон: {tone}\n\n"
            f"Язык письма: {lang}"
        )

    if mode == "concept":
        brief = data.get("brief", "")
        dmc_program = data.get("dmc_program", "").strip()
        msg = f"Бриф:\n{brief}"
        if dmc_program:
            msg += f"\n\nПрограмма от DMC:\n{dmc_program}"
        else:
            msg += "\n\nПрограмма от DMC: не предоставлена (используй Сценарий А)"
        return msg

    return ""


PROMPT_MAP = {
    "brief": p.BRIEF_PROMPT,
    "objection": p.OBJECTION_PROMPT,
    "supplier": p.SUPPLIER_PROMPT,
    "dmc": p.DMC_PROMPT,
    "dmc_refine": p.DMC_REFINE_PROMPT,
    "followup": p.FOLLOWUP_PROMPT,
    "concept": p.CONCEPT_PROMPT,
}


def _extract_pdf(content):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    parts = [page.extract_text() for page in reader.pages if page.extract_text()]
    return "\n\n".join(parts)


def _extract_docx(content):
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(content):
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[{sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_pptx(content):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        texts.append(line)
        if texts:
            slides.append(f"[Слайд {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


@app.route("/parse-file", methods=["POST"])
def parse_file():
    if "file" not in request.files:
        return jsonify({"error": "Файл не найден"}), 400
    file = request.files["file"]
    if not file.filename:
        return jsonify({"error": "Файл не выбран"}), 400

    filename = file.filename.lower()
    content = file.read()

    try:
        if filename.endswith(".pdf"):
            text = _extract_pdf(content)
        elif filename.endswith(".docx"):
            text = _extract_docx(content)
        elif filename.endswith(".xlsx"):
            text = _extract_xlsx(content)
        elif filename.endswith(".pptx"):
            text = _extract_pptx(content)
        elif filename.endswith(".xls"):
            return jsonify({"error": "Формат .xls не поддерживается. Сохраните файл как .xlsx (Excel 2007+)."}), 400
        elif filename.endswith(".doc"):
            return jsonify({"error": "Формат .doc не поддерживается. Сохраните файл как .docx."}), 400
        elif filename.endswith(".ppt"):
            return jsonify({"error": "Формат .ppt не поддерживается. Сохраните файл как .pptx (PowerPoint 2007+)."}), 400
        else:
            return jsonify({"error": "Неподдерживаемый формат. Используйте PDF, Word (.docx), Excel (.xlsx) или PowerPoint (.pptx)."}), 400

        if not text.strip():
            return jsonify({"error": "Не удалось извлечь текст из файла. Возможно, файл содержит только изображения."}), 400

        return jsonify({"text": text})
    except Exception as e:
        return jsonify({"error": f"Ошибка при чтении файла: {str(e)}"}), 500


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    data = request.json
    mode = data.get("mode")
    system_prompt = PROMPT_MAP.get(mode)

    if not system_prompt:
        return Response(
            f"data: {json.dumps({'error': 'Неизвестный режим'})}\n\n",
            content_type="text/event-stream",
        )

    user_message = build_user_message(mode, data)

    def stream():
        try:
            with client.messages.stream(
                model=MODEL,
                max_tokens=MAX_TOKENS,
                system=system_prompt,
                messages=[{"role": "user", "content": user_message}],
            ) as s:
                for text in s.text_stream:
                    yield f"data: {json.dumps({'chunk': text})}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return Response(
        stream_with_context(stream()),
        content_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


if __name__ == "__main__":
    app.run(debug=True, port=8080)
